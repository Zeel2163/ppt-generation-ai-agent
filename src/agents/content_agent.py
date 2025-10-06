from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

class ContentAgent:
    def __init__(self, llm_service, image_service=None):
        self.llm_service = llm_service
        self.image_service = image_service
        self.presentation = Presentation()

    def create_title_slide(self, title, subtitle=""):
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(30)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(20)
            subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def create_content_slide(self, title, content, image_query=None):
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        body_shape = slide.placeholders[1]
        body_shape.text = content
        for paragraph in body_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(20)

        if image_query and self.image_service:
            image_path = self.image_service.download_image(image_query)
            if image_path and os.path.exists(image_path):
                slide.shapes.add_picture(image_path, Inches(6), Inches(2), height=Inches(4))
                os.remove(image_path)

    def generate_presentation(self, topic, num_slides=5, output_file="presentation.pptx"):
        outline = self.llm_service.generate_content_outline(topic, num_slides)
        if not outline:
            print("Failed to generate content outline.")
            return None

        for i, slide_data in enumerate(outline):
            title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            slide_type = slide_data.get("slide_type", "content")
            image_query = slide_data.get("image_query")

            if slide_type == "title":
                self.create_title_slide(title, "Created by AI Agent")
            elif slide_type == "image":
                self.create_content_slide(title, content, image_query=image_query)
            else:
                self.create_content_slide(title, content, image_query=(image_query if i % 3 == 0 else None))

        self.presentation.save(output_file)
        print(f"Presentation saved to {output_file}")
        return output_file
