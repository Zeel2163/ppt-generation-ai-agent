import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from src.services.llm_service import LLMService
from src.services.image_service import ImageService
from src.config.settings import PEXELS_API_KEY


class PPTGenerator:
    def __init__(self, api_key: str):
        self.llm = LLMService(api_key)
        self.image_service = ImageService(PEXELS_API_KEY)
        self.presentation = Presentation()

    def create_title_slide(self, title, subtitle=""):
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(30)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(20)
            subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    from pptx.util import Pt, Inches

    def create_content_slide(self, title, content, include_image=False):
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()

        try:
            if isinstance(content, list):
                # First bullet
                p = text_frame.paragraphs[0]
                p.text = content[0]
                p.font.size = Pt(20)

                # Add remaining as new bullets
                for item in content[1:]:
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.level = 0
                    p.font.size = Pt(20)
            else:
                # If content is just a string
                text_frame.text = content
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(20)

        except Exception as e:
            print(f"getting exception in content shaping: {e}")

        if include_image:
            image_desc = self.llm.generate_image_description(
                " ".join(content) if isinstance(content, list) else content
            )
            try:
                image_path = self.image_service.download_image(image_desc)
                slide.shapes.add_picture(image_path, Inches(6), Inches(2), height=Inches(4))
                os.remove(image_path)
            except Exception as e:
                print(f"Image error in create_content_slide: {e}")

    # def create_content_slide(self, title, content, include_image=False):
    #     slide_layout = self.presentation.slide_layouts[1]
    #     slide = self.presentation.slides.add_slide(slide_layout)
    #
    #     slide.shapes.title.text = title
    #     content_shape = slide.placeholders[1]
    #     try:
    #         content_shape.text = content
    #     except Exception as e:
    #         print(f"getting exception in content shaping: {e}")
    #
    #     for paragraph in content_shape.text_frame.paragraphs:
    #         paragraph.font.size = Pt(20)
    #
    #     if include_image:
    #         image_desc = self.llm.generate_image_description(content)
    #         try:
    #             image_path = self.image_service.download_image(image_desc)
    #             slide.shapes.add_picture(image_path, Inches(6), Inches(2), height=Inches(4))
    #             os.remove(image_path)
    #         except Exception as e:
    #             print(f"Image error in create_content_slide: {e}")

    def generate_presentation(self, topic, num_slides=5, output_file="presentation.pptx"):
        content_outline = self.llm.generate_content_outline(topic, num_slides)

        for i, slide_data in enumerate(content_outline):
            title = slide_data["title"]
            content = slide_data["content"]
            slide_type = slide_data["slide_type"]

            if slide_type == "title":
                self.create_title_slide(title, "Created by Zeel")

            elif slide_type == "image":
                # img_query = self.llm.generate_image_description(content)
                self.create_content_slide(title, content, include_image=True)

            else:
                include_image = (i % 3 == 0)
                self.create_content_slide(title, content, include_image)

        self.presentation.save(output_file)
        return output_file
